"""
AURA NX-Alpha — Document Parser
Extract plain text from uploaded PDF and DOCX files.
Returns the full extracted text string; chunking for memory storage is
handled downstream by data_controller._ingest_personal_document().
"""

from __future__ import annotations

import io
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".csv", ".json", ".xml", ".yaml", ".yml",
    ".toml", ".ini", ".cfg", ".env", ".log", ".sql", ".rtf",
    ".py", ".js", ".ts", ".jsx", ".tsx", ".html", ".htm", ".css", ".scss",
    ".sass", ".less", ".sh", ".bash", ".ps1", ".rb", ".go", ".rs", ".java",
    ".cpp", ".c", ".h", ".php", ".lua", ".swift", ".kt", ".r", ".pl",
    ".ex", ".hs", ".elm", ".ml", ".clj", ".scala", ".zig", ".tf",
}


def extract_text(file_bytes: bytes, filename: str) -> str:
    """
    Extract plain text from a document file.

    Supported formats:
      - .pdf       — via pypdf
      - .docx/.doc — via python-docx
      - .xlsx      — via openpyxl (all sheets, cell values as text)
      - .pptx      — via python-pptx (all slide text)
      - text/code  — UTF-8 decode (see TEXT_EXTENSIONS)

    Args:
        file_bytes: Raw file content.
        filename:   Original filename (used to detect format via extension).

    Returns:
        Extracted text as a single string.

    Raises:
        ValueError: Unsupported file format.
        RuntimeError: Extraction failed (corrupt file, missing dependency, etc.).
    """
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(file_bytes, filename)
    elif ext in (".docx", ".doc"):
        return _extract_docx(file_bytes, filename)
    elif ext in (".xlsx", ".xls"):
        return _extract_xlsx(file_bytes, filename)
    elif ext == ".pptx":
        return _extract_pptx(file_bytes, filename)
    elif ext in TEXT_EXTENSIONS:
        return _extract_text_file(file_bytes, filename)
    else:
        raise ValueError(
            f"Unsupported format '{ext}'. "
            "Supported: .pdf, .docx, .xlsx, .pptx, and plain text/code files."
        )


def _extract_pdf(data: bytes, filename: str) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as e:
        raise RuntimeError(
            "pypdf is not installed. Run: pip install pypdf>=4.0"
        ) from e

    try:
        reader = PdfReader(io.BytesIO(data))
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(text)
        full_text = "\n\n".join(pages).strip()
        logger.info("[document_parser] PDF '%s': %d pages, %d chars", filename, len(reader.pages), len(full_text))
        if not full_text:
            raise RuntimeError(
                f"No text extracted from '{filename}'. "
                "The PDF may be image-only (scanned). Vision-based OCR is not yet supported."
            )
        return full_text
    except Exception as exc:
        if "pypdf" not in str(type(exc).__module__):
            raise
        raise RuntimeError(f"Failed to parse PDF '{filename}': {exc}") from exc


def _extract_docx(data: bytes, filename: str) -> str:
    try:
        from docx import Document
    except ImportError as e:
        raise RuntimeError(
            "python-docx is not installed. Run: pip install python-docx"
        ) from e

    try:
        doc = Document(io.BytesIO(data))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        paragraphs.append(cell.text.strip())
        full_text = "\n\n".join(paragraphs)
        logger.info("[document_parser] DOCX '%s': %d paragraphs, %d chars", filename, len(paragraphs), len(full_text))
        if not full_text:
            raise RuntimeError(f"No text content found in '{filename}'.")
        return full_text
    except Exception as exc:
        raise RuntimeError(f"Failed to parse DOCX '{filename}': {exc}") from exc


def _extract_xlsx(data: bytes, filename: str) -> str:
    try:
        import openpyxl
    except ImportError as e:
        raise RuntimeError("openpyxl is not installed. Run: pip install openpyxl") from e

    try:
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        sections = []
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows.append("\t".join(cells))
            if rows:
                sections.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
        wb.close()
        full_text = "\n\n".join(sections)
        logger.info("[document_parser] XLSX '%s': %d sheet(s), %d chars", filename, len(sections), len(full_text))
        if not full_text:
            raise RuntimeError(f"No cell data found in '{filename}'.")
        return full_text
    except Exception as exc:
        raise RuntimeError(f"Failed to parse XLSX '{filename}': {exc}") from exc


def _extract_pptx(data: bytes, filename: str) -> str:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx") from e

    try:
        prs = Presentation(io.BytesIO(data))
        slides = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs).strip()
                        if line:
                            texts.append(line)
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        full_text = "\n\n".join(slides)
        logger.info("[document_parser] PPTX '%s': %d slide(s), %d chars", filename, len(slides), len(full_text))
        if not full_text:
            raise RuntimeError(f"No text content found in '{filename}'.")
        return full_text
    except Exception as exc:
        raise RuntimeError(f"Failed to parse PPTX '{filename}': {exc}") from exc


def _extract_text_file(data: bytes, filename: str) -> str:
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        try:
            text = data.decode("latin-1")
        except Exception as e:
            raise RuntimeError(f"Could not decode '{filename}' as text: {e}") from e
    text = text.strip()
    if not text:
        raise RuntimeError(f"No text content found in '{filename}'.")
    logger.info("[document_parser] Text file '%s': %d chars", filename, len(text))
    return text
